import subprocess
try:
    from pptx import Presentation
except:
    subprocess.run(['pip', 'install', 'python-pptx'], capture_output=True)
    from pptx import Presentation

import os
outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ch8_pptx_output.txt')
with open(outpath, 'w', encoding='utf-8') as out:
    for fname in ['gw/chapter_wise/Chapter 8.pptx', 'gw/chapter_wise/Chapter 8_2.pptx']:
        out.write(f'\n====== {fname} ======\n')
        prs = Presentation(fname)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                out.write(f'--- SLIDE {i+1} ---\n')
                for t in texts:
                    out.write(t + '\n')
print('Done - output written to ' + outpath)
