import fitz
from pptx import Presentation
import sys

outpath = 'gw/ch9_ch6_output.txt'
with open(outpath, 'w', encoding='utf-8') as f:
    # PDF
    f.write('====== Chapter 9.pdf ======\n')
    doc = fitz.open('gw/chapter_wise/Chapter 9.pdf')
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            f.write(f'--- PAGE {i+1} ---\n')
            f.write(text + '\n')

    # PPTX
    f.write('\n====== 6.GW.pptx ======\n')
    prs = Presentation('gw/chapter_wise/6.GW.pptx')
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            f.write(f'--- SLIDE {i+1} ---\n')
            for t in texts:
                f.write(t + '\n')

print(f'Done. Output written to {outpath}')
